from pptx import Presentation

# Create a presentation object
ppt = Presentation()

# Slide 1: Title Slide
slide_1 = ppt.slides.add_slide(ppt.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Statistics I: Mean, Median, Mode, Variance"
subtitle.text = "Understanding Key Concepts in Data Analysis\nYour Name\nDate"

# Slide 2: Introduction to Statistics
slide_2 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_2.shapes.title
content = slide_2.shapes.placeholders[1]
title.text = "Introduction to Statistics"
content.text = "Statistics is essential in data science for summarizing and interpreting data.\n\nReal-world applications: surveys, market research, scientific studies."

# Slide 3: What is Mean?
slide_3 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_3.shapes.title
content = slide_3.shapes.placeholders[1]
title.text = "What is Mean?"
content.text = "Definition: The average of a set of numbers\nFormula: Mean = Σx / n\nExample: Find the mean of [5, 10, 15, 20]"

# Slide 4: Live Example - Calculating Mean
slide_4 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_4.shapes.title
content = slide_4.shapes.placeholders[1]
title.text = "Live Example: Calculating Mean"
content.text = "Example Calculation: (5 + 10 + 15 + 20) / 4 = 12.5"

# Slide 5: What is Median?
slide_5 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_5.shapes.title
content = slide_5.shapes.placeholders[1]
title.text = "What is Median?"
content.text = "Definition: The middle value when data is sorted.\nExample: Find the median of [3, 1, 4, 5, 2]"

# Slide 6: Live Example - Finding the Median
slide_6 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_6.shapes.title
content = slide_6.shapes.placeholders[1]
title.text = "Live Example: Finding the Median"
content.text = "Steps:\n1. Sort data: [1, 2, 3, 4, 5]\n2. The middle value is 3."

# Slide 7: What is Mode?
slide_7 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_7.shapes.title
content = slide_7.shapes.placeholders[1]
title.text = "What is Mode?"
content.text = "Definition: The value that appears most frequently.\nExample: Find the mode of [1, 2, 2, 3, 3, 3, 4]"

# Slide 8: Live Example - Finding the Mode
slide_8 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_8.shapes.title
content = slide_8.shapes.placeholders[1]
title.text = "Live Example: Finding the Mode"
content.text = "Mode = 3 (appears 3 times)\nThere is one mode in this case."

# Slide 9: What is Variance?
slide_9 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_9.shapes.title
content = slide_9.shapes.placeholders[1]
title.text = "What is Variance?"
content.text = "Definition: Measures the spread of data points from the mean.\nFormula: Variance = Σ(x - μ)² / n\nExample: Calculate variance for [5, 10, 15, 20]."

# Slide 10: Live Example - Calculating Variance
slide_10 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_10.shapes.title
content = slide_10.shapes.placeholders[1]
title.text = "Live Example: Calculating Variance"
content.text = "Steps:\n1. Calculate the mean: 12.5\n2. Calculate squared differences: [(5-12.5)², (10-12.5)², (15-12.5)², (20-12.5)²]\n3. Find the average of squared differences."

# Slide 11: Why Do We Need These Measures?
slide_11 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_11.shapes.title
content = slide_11.shapes.placeholders[1]
title.text = "Why Do We Need These Measures?"
content.text = "Mean: Central tendency of the data\nMedian: Useful for skewed data\nMode: Most frequent data points\nVariance: Spread and consistency of data"

# Slide 12: Real-World Application: Comparing Data
slide_12 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_12.shapes.title
content = slide_12.shapes.placeholders[1]
title.text = "Real-World Application: Comparing Data"
content.text = "Example: Comparing exam scores from two classes and understanding their performance using statistics."

# Slide 13: Key Differences Between Mean, Median, and Mode
slide_13 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_13.shapes.title
content = slide_13.shapes.placeholders[1]
title.text = "Key Differences Between Mean, Median, and Mode"
content.text = "Mean: Sensitive to outliers\nMedian: Robust in the presence of outliers\nMode: Helps identify common patterns"

# Slide 14: Conclusion
slide_14 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_14.shapes.title
content = slide_14.shapes.placeholders[1]
title.text = "Conclusion"
content.text = "Recap of the key points: Mean, Median, Mode, and Variance.\nPractice with real datasets to deepen understanding."

# Slide 15: Q&A and Wrap-Up
slide_15 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_15.shapes.title
content = slide_15.shapes.placeholders[1]
title.text = "Q&A and Wrap-Up"
content.text = "Open the floor for questions.\nProvide further resources for learning."

# Save the presentation to a file
ppt_filename = "Statistics_Introduction.pptx"
ppt.save(ppt_filename)


